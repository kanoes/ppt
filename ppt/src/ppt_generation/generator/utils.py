"""Utility helpers supporting PPT generation."""

import copy
import io
import os
import random
import re
import time

import pptx_ea_font
from dotenv import load_dotenv
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from shared.logging import get_logger
from shared.services.llm import LLM

load_dotenv()

# ログ設定
logger = get_logger("ppt_utils")


class LLMInvoker:
    def __init__(self, deployment_name: str | None = None, temperature: float | None = None, json_mode: bool = False):
        default_dep = os.getenv("DEFAULT_LLM_DEPLOYMENT", "gpt-5")
        default_temp = float(os.getenv("DEFAULT_LLM_TEMPERATURE", "0.3"))

        self.deployment = deployment_name or default_dep
        self.temperature = default_temp if temperature is None else float(temperature)

        self.llm = LLM(
            deployment_name=self.deployment,
            temperature=self.temperature,
            json_mode=json_mode,
        )

    def invoke(self, prompt_template: str, **kwargs) -> str:
        """
        テンプレート文字列を format(**kwargs) して LLM を叩き、文字列を返す。
        """
        try:
            prompt_text = prompt_template.format(**kwargs)
        except KeyError as e:
            logger.error({
                "message": "プロンプトテンプレートに必要な引数が不足しています。",
                "missing_argument": str(e),
                "status": "problem",
            })
            raise

        try:
            logger.info({
                "message": "LLM呼び出し開始",
                "operation": "ppt_llm_invoke",
                "deployment": self.deployment,
                "temperature": self.temperature,
                "status": "started",
            })
            start_time = time.time()

            answer = self.llm.invoke(prompt_text)

            end_time = time.time()
            execution_time = end_time - start_time

            usage_new = getattr(answer, "usage_metadata", None) or {}
            resp_meta = getattr(answer, "response_metadata", {}) or {}
            usage_old = resp_meta.get("token_usage", {}) if isinstance(resp_meta, dict) else {}

            token_log = {
                "input_tokens": usage_new.get("input_tokens"),
                "output_tokens": usage_new.get("output_tokens"),
                "total_tokens": usage_new.get("total_tokens") or usage_old.get("total_tokens"),
                "prompt_tokens": usage_old.get("prompt_tokens"),
                "completion_tokens": usage_old.get("completion_tokens"),
                "model": resp_meta.get("model") if isinstance(resp_meta, dict) else None,
                "system_fingerprint": resp_meta.get("system_fingerprint") if isinstance(resp_meta, dict) else None,
                "deployment_name": self.deployment,
                "temperature": self.temperature,
                "execution_time": execution_time,
            }
            logger.info({
                "message": "LLMトークン使用量",
                "operation": "llm_invoke_usage",
                "tokens": token_log,
            })

            content = getattr(answer, "content", answer)
            if not isinstance(content, str):
                raise TypeError("応答が文字列ではありません。")

            logger.info({
                "message": "LLM呼び出し完了",
                "operation": "ppt_llm_invoke",
                "status": "completed",
            })
            return content

        except Exception as e:
            logger.error({
                "message": "LLMの呼び出し中にエラーが発生しました。",
                "operation": "ppt_llm_invoke",
                "error_message": str(e),
                "status": "problem",
            })
            raise


class PPTUtils:
    """
    PowerPoint操作に関するユーティリティクラス。
    """

    @staticmethod
    def duplicate_slide(index, presentation):
        """
        指定されたインデックスのスライドを複製する。
        """
        try:
            template_slide = presentation.slides[index]
            copied_slide = presentation.slides.add_slide(template_slide.slide_layout)

            # 新しいスライドをクリア
            for shape in list(copied_slide.shapes):
                copied_slide.shapes.element.remove(shape.element)

            # 元スライドの全てのシェイプを複製
            for shape in template_slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = io.BytesIO(shape.image.blob)
                    copied_slide.shapes.add_picture(
                        image_file=img,
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                    )
                else:
                    new_element = copy.deepcopy(shape.element)
                    copied_slide.shapes._spTree.insert_element_before(
                        new_element, "p:extLst"
                    )

            return copied_slide

        except Exception as e:
            logger.error({
                "message": f"スライド {index} の複製中にエラーが発生しました。",
                "error_message": str(e)
            })
            raise

    @staticmethod
    def remove_original_slides(presentation, original_slide_count):
        """
        指定された数の元スライドをプレゼンテーションから削除する。
        """
        try:
            xml_slides = presentation.slides._sldIdLst
            slide_ids = list(xml_slides)[:original_slide_count]
            for slide_id in slide_ids:
                xml_slides.remove(slide_id)
        except Exception as e:
            logger.error({
                "message": "元スライドの削除中にエラーが発生しました。",
                "error_message": str(e)
            })
            raise

    @staticmethod
    def add_text_to_shape(shape, text, hyperlink=None):
        """
        シェイプにテキストを追加し、フォントスタイルを継承しながらハイパーリンクを追加する。
        """
        try:
            text_frame = shape.text_frame
            
            # <br>タグを実際の改行に置換する
            text = text.replace('<br>', '\n')

            # 既存のフォント情報を取得する
            if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                first_run = text_frame.paragraphs[0].runs[0]
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_color = (
                    font.color.rgb if font.color.type == MSO_COLOR_TYPE.RGB else None
                )
            else:
                # デフォルトのフォント設定
                font_name = "Meiryo UI"
                font_size = Pt(18)
                font_bold = False
                font_italic = False
                font_color = RGBColor(0, 0, 0)

            # テキストをクリアし、新しいテキストを設定する
            text_frame.clear()
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = text
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = font_bold
            run.font.italic = font_italic
            if font_color:
                run.font.color.rgb = font_color
            pptx_ea_font.set_font(run, 'Meiryo UI')
            run.font.name = font_name

            # ハイパーリンクを追加（オプション）
            if hyperlink:
                run.hyperlink.address = hyperlink
                run.font.color.rgb = RGBColor(0, 0, 255)
                run.font.underline = True
                
        except Exception as e:
            logger.error({
                "message": "シェイプにテキストを追加中にエラーが発生しました。",
                "error_message": str(e)
            })
            raise


    @staticmethod
    def add_picture_to_slide(placeholder, chart):
        """
        指定されたスライドに画像を挿入する。
        """
        try:
            placeholder.insert_picture(chart)
        except Exception as e:
            logger.error({
                "message": "スライドへの画像挿入中にエラーが発生しました。",
                "error_message": str(e)
            })
            raise

    @staticmethod
    def extract_all_between_tags(tag, text):
        """
        指定されたタグに囲まれたテキストを全て抽出する。
        """
        try:
            start_tag = f"[{tag}]"
            end_tag = f"[/{tag}]"
            pattern = re.escape(start_tag) + r"(.*?)" + re.escape(end_tag)
            results = [match.strip() for match in re.findall(pattern, (text or ""), re.DOTALL)]
            return results
        except Exception as e:
            logger.error({
                "message": f"タグ '{tag}' の抽出中にエラーが発生しました。",
                "error_message": str(e)
            })
            raise



class OtherUtils:
    _previous_templates = []

    @staticmethod
    def random_choice(template_info_group):
        """
        テンプレート情報のグループからランダムに選択する。
        """
        try:
            if (
                len(OtherUtils._previous_templates) >= 2
                and OtherUtils._previous_templates[-1] == OtherUtils._previous_templates[-2]
            ):
                # 直近の重複を避ける
                variant_key = [x for x in template_info_group if x != OtherUtils._previous_templates[-1]][0]
            else:
                variant_key = random.choice(template_info_group)

            OtherUtils._previous_templates.append(variant_key)
            if len(OtherUtils._previous_templates) > 2:
                OtherUtils._previous_templates.pop(0)

            return variant_key

        except Exception as e:
            logger.error({
                "message": "ランダム選択機能にエラーが発生しました。",
                "error_message": str(e)
            })
            raise
