"""Manual PPT manipulation test helpers."""

import copy
import io
import json
import re
from pathlib import Path

import pptx_ea_font
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

PROJECT_ROOT = Path(__file__).resolve().parents[2]
RESOURCES_DIR = PROJECT_ROOT / "ppt" / "resources"
REQUIRED_FILES = [
    "test_title_template_info.json",
    "test_chart_template_info.json",
    "test_reference_template_info.json",
    "test_normal_template_info.json",
]

MISSING_FILES = [f for f in REQUIRED_FILES if not (RESOURCES_DIR / f).exists()]
if MISSING_FILES:
    pytest.skip(
        f"Skipping template insertion tests; missing resources: {', '.join(MISSING_FILES)}",
        allow_module_level=True,
    )


# テンプレート情報の読み込み
with open(RESOURCES_DIR / "test_title_template_info.json", "r") as file:
    title_template_info = json.load(file)

with open(RESOURCES_DIR / "test_chart_template_info.json", "r") as file:
    chart_template_info = json.load(file)

with open(RESOURCES_DIR / "test_reference_template_info.json", "r") as file:
    reference_template_info = json.load(file)

with open(RESOURCES_DIR / "test_normal_template_info.json", "r") as file:
    normal_template_info = json.load(file)
    
    
def duplicate_slide(index, presentation):
    """
    指定されたインデックスのスライドを複製する。
    """
    try:
        template_slide = presentation.slides[index]
        copied_slide = presentation.slides.add_slide(template_slide.slide_layout)

        # 新しいスライドをクリア
        for shape in list(copied_slide.shapes):
            copied_slide.shapes.element.remove(shape.element)

        # 元スライドの全てのシェイプを複製
        for shape in template_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = io.BytesIO(shape.image.blob)
                copied_slide.shapes.add_picture(
                    image_file=img,
                    left=shape.left,
                    top=shape.top,
                    width=shape.width,
                    height=shape.height,
                )
            else:
                new_element = copy.deepcopy(shape.element)
                copied_slide.shapes._spTree.insert_element_before(
                    new_element, "p:extLst"
                )

        return copied_slide
    except Exception as e:
        raise
    
def add_text_to_shape(shape, text, hyperlink=None):
        """
        シェイプにテキストを追加し、フォントスタイルを継承しながらハイパーリンクを追加する。
        """
        text_frame = shape.text_frame

        # 既存のフォント情報を取得する
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            first_run = text_frame.paragraphs[0].runs[0]
            font = first_run.font
            font_name = font.name
            print(font_name)
            font_size = font.size
            font_bold = font.bold
            font_italic = font.italic
            font_color = (
                font.color.rgb if font.color.type == MSO_COLOR_TYPE.RGB else None
            )
        else:
            # デフォルトのフォント設定
            font_name = "Meiryo UI"
            font_size = Pt(18)
            font_bold = False
            font_italic = False
            font_color = RGBColor(0, 0, 0)

        # テキストをクリアし、新しいテキストを設定する
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
        run.font.name = font_name
        pptx_ea_font.set_font(run, 'Meiryo UI')
        run.font.size = font_size
        run.font.bold = font_bold
        run.font.italic = font_italic
        if font_color:
            run.font.color.rgb = font_color

        # ハイパーリンクを追加（オプション）
        if hyperlink:
            run.hyperlink.address = hyperlink
            run.font.color.rgb = RGBColor(0, 0, 255)
            run.font.underline = True

def add_picture_to_slide(placeholder, chart):
    """
    指定されたスライドに画像を挿入する。
    """
    try:
        # picture = placeholder.insert_picture(chart)
        placeholder.insert_picture(chart)
    except Exception as e:
        raise

@staticmethod
def extract_all_between_tags(tag, text):
    """
    指定されたタグに囲まれたテキストを全て抽出する。
    """
    try:
        start_tag = f"[{tag}]"
        end_tag = f"[/{tag}]"
        pattern = re.escape(start_tag) + r"(.*?)" + re.escape(end_tag)
        results = [match.strip() for match in re.findall(pattern, text, re.DOTALL)]
        return results
    except Exception as e:
        raise

def test_insert_content():
    template_path = RESOURCES_DIR / "smbc_template_new.pptx"
    presentation = Presentation(template_path)

    # 目標テンプレートの情報を取得
    template_info_standard = "1"  # ←←←標準テンプレートID
    template_info_item = next(
        (
            item["mainA"]
            for item in normal_template_info
            if item["template_id"] == template_info_standard
        ),
        None,
    )

    # 目標シェイプや目標プレースホルダーの番号を取得
    title_placeholder_number = template_info_item["placeholder_number"]["title"]
    subtitle_shape_number = template_info_item["shape_number"]["subtitle"]
    body_shape_number = template_info_item["shape_number"]["body"]

    # 目標テンプレートスライドを複製
    slide_number = template_info_item["slide_number"]
    slide = duplicate_slide(slide_number, presentation)

    # コンテンツを埋め込む
    title = "Sample Title"
    subtitle = "Sample Subtitle"
    add_text_to_shape(slide.placeholders[title_placeholder_number], "例文")
    add_text_to_shape(slide.shapes[subtitle_shape_number], "例文")
    add_text_to_shape(slide.shapes[body_shape_number], "テスト文だけですテスト文だけですテスト文だけですテスト文だけですテスト文だけですテスト文だけですテスト文だけですテスト文だけですテスト文だけですテスト文だけですテスト文だけですテスト文だけですテスト文だけですテスト文だけです")
    presentation.save(template_path)


if __name__ == "__main__":
    test_insert_content()
    