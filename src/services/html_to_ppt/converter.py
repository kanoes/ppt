from __future__ import annotations

import base64
from dataclasses import dataclass, field
from html.parser import HTMLParser
from io import BytesIO
from typing import List, Sequence, Tuple

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.logging import get_logger

logger = get_logger("html_to_ppt")


@dataclass
class SlideData:
    role: str
    headings: list[Tuple[str, str]] = field(default_factory=list)
    paragraphs: list[Tuple[str, str]] = field(default_factory=list)
    lists: list[list[str]] = field(default_factory=list)
    images: list[Tuple[str, str | None]] = field(default_factory=list)


class _SlideHTMLParser(HTMLParser):
    """Very small HTML parser tailored to DeckHTMLRenderer output."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.slides: list[SlideData] = []
        self._current_slide: SlideData | None = None
        self._tag_stack: list[str] = []
        self._attr_stack: list[dict[str, str]] = []
        self._text_stack: list[str | None] = []
        self._list_stack: list[list[str]] = []
        self._figure_stack: list[dict[str, str | None]] = []

    def handle_starttag(self, tag: str, attrs: list[Tuple[str, str | None]]) -> None:
        attr_dict = {name: value or "" for name, value in attrs}
        self._tag_stack.append(tag)
        self._attr_stack.append(attr_dict)

        if tag == "section":
            classes = attr_dict.get("class", "").split()
            if "slide" in classes:
                role = attr_dict.get("data-role", "content").lower()
                slide = SlideData(role=role)
                self.slides.append(slide)
                self._current_slide = slide
            self._text_stack.append(None)
            return

        if self._current_slide is None:
            self._text_stack.append(None)
            return

        if tag in {"h1", "h2", "h3", "p", "figcaption", "li"}:
            self._text_stack.append("")
        else:
            self._text_stack.append(None)

        if tag in {"ul", "ol"}:
            self._list_stack.append([])
        elif tag == "figure":
            self._figure_stack.append({"src": None, "caption": None})
        elif tag == "img":
            src = attr_dict.get("src")
            if self._figure_stack:
                self._figure_stack[-1]["src"] = src
            elif self._current_slide and src:
                self._current_slide.images.append((src, None))

    def handle_endtag(self, tag: str) -> None:
        if not self._tag_stack:
            return
        current_tag = self._tag_stack.pop()
        attrs = self._attr_stack.pop()
        text = self._text_stack.pop() if self._text_stack else None

        if current_tag != tag:
            # mismatched tag – ignore for robustness
            return

        if tag == "section":
            self._current_slide = None
            return

        if self._current_slide is None:
            return

        stripped = text.strip() if isinstance(text, str) else ""

        if tag in {"h1", "h2", "h3"}:
            if stripped:
                self._current_slide.headings.append((tag, stripped))
        elif tag == "p":
            if stripped:
                class_name = attrs.get("class", "")
                self._current_slide.paragraphs.append((class_name, stripped))
        elif tag == "li":
            if stripped and self._list_stack:
                self._list_stack[-1].append(stripped)
        elif tag in {"ul", "ol"}:
            if self._list_stack:
                items = self._list_stack.pop()
                if items:
                    self._current_slide.lists.append(items)
        elif tag == "figcaption":
            if stripped and self._figure_stack:
                self._figure_stack[-1]["caption"] = stripped
        elif tag == "figure":
            if self._figure_stack:
                fig = self._figure_stack.pop()
                src = fig.get("src")
                if src:
                    self._current_slide.images.append((src, fig.get("caption")))

    def handle_data(self, data: str) -> None:
        if not self._text_stack:
            return
        if self._text_stack[-1] is None:
            return
        self._text_stack[-1] += data


TextBlock = Tuple[str, Sequence[str]]
ImageInfo = Tuple[str, str | None]


class HTMLToPPTXConverter:
    """Convert structured HTML slides into a PowerPoint presentation."""

    def __init__(self, width_in: float = 10.0, height_in: float = 5.625) -> None:
        self.width_in = width_in
        self.height_in = height_in

    def convert(self, html_text: str) -> BytesIO:
        parser = _SlideHTMLParser()
        parser.feed(html_text or "")

        presentation = Presentation()
        presentation.slide_width = Inches(self.width_in)
        presentation.slide_height = Inches(self.height_in)

        for slide_data in parser.slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            role = slide_data.role
            if role == "title":
                self._populate_title_slide(slide, slide_data)
            elif role == "chart":
                self._populate_chart_slide(slide, slide_data)
            elif role == "sources":
                self._populate_sources_slide(slide, slide_data)
            else:
                self._populate_content_slide(slide, slide_data)

        buffer = BytesIO()
        presentation.save(buffer)
        buffer.seek(0)
        return buffer

    # ------------------------------------------------------------------
    # Slide population helpers
    # ------------------------------------------------------------------
    def _populate_title_slide(self, slide, slide_data: SlideData) -> None:
        title = self._select_heading(slide_data) or "プレゼンテーション"
        subtitle = self._find_paragraph(slide_data, "subtitle")
        eyebrow = self._find_paragraph(slide_data, "eyebrow")

        self._add_title_box(slide, title)
        paragraphs: list[TextBlock] = []
        if eyebrow:
            paragraphs.append(("paragraph", [eyebrow]))
        if subtitle:
            paragraphs.append(("paragraph", [subtitle]))
        if paragraphs:
            self._add_text_blocks(slide, paragraphs, top_inches=2.0)

    def _populate_content_slide(self, slide, slide_data: SlideData) -> None:
        heading = self._select_heading(slide_data)
        if heading:
            self._add_title_box(slide, heading)
        lead = self._find_paragraph(slide_data, "lead")
        body_paragraphs = list(
            self._iter_paragraphs(slide_data, exclude_classes={"lead", "subtitle", "eyebrow"})
        )
        text_blocks: list[TextBlock] = []
        if lead:
            text_blocks.append(("paragraph", [lead]))
        for paragraph in body_paragraphs:
            text_blocks.append(("paragraph", [paragraph]))
        for items in slide_data.lists:
            text_blocks.append(("list", items))
        if text_blocks:
            self._add_text_blocks(slide, text_blocks)
        if slide_data.images:
            self._add_images(slide, slide_data.images, top_inches=3.4 if text_blocks else 1.8)

    def _populate_chart_slide(self, slide, slide_data: SlideData) -> None:
        heading = self._select_heading(slide_data)
        if heading:
            self._add_title_box(slide, heading)
        lead = self._find_paragraph(slide_data, "lead")
        if lead:
            self._add_text_blocks(slide, [("paragraph", [lead])])
        if slide_data.images:
            self._add_images(slide, slide_data.images, top_inches=1.8)

    def _populate_sources_slide(self, slide, slide_data: SlideData) -> None:
        heading = self._select_heading(slide_data) or "引用元"
        self._add_title_box(slide, heading)
        items: list[str] = []
        if slide_data.lists:
            for lst in slide_data.lists:
                items.extend(lst)
        else:
            items.extend(self._iter_paragraphs(slide_data))
        if not items:
            items.append("引用元情報は提供されませんでした。")
        self._add_text_blocks(slide, [("list", items)])

    # ------------------------------------------------------------------
    # Shape helpers
    # ------------------------------------------------------------------
    def _add_title_box(self, slide, text: str) -> None:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(1.2))
        frame = box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text.strip()
        paragraph.font.size = Pt(34)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.LEFT

    def _add_text_blocks(
        self,
        slide,
        blocks: Sequence[TextBlock],
        *,
        top_inches: float = 1.8,
    ) -> None:
        if not blocks:
            return
        box = slide.shapes.add_textbox(Inches(0.9), Inches(top_inches), Inches(8.2), Inches(3.8))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True

        first_paragraph = True
        for block_type, values in blocks:
            if block_type == "paragraph":
                for value in values:
                    paragraph = frame.paragraphs[0] if first_paragraph else frame.add_paragraph()
                    paragraph.text = value
                    paragraph.font.size = Pt(20)
                    paragraph.level = 0
                    paragraph.alignment = PP_ALIGN.LEFT
                    first_paragraph = False
            elif block_type == "list":
                for value in values:
                    paragraph = frame.paragraphs[0] if first_paragraph else frame.add_paragraph()
                    paragraph.text = value
                    paragraph.font.size = Pt(18)
                    paragraph.level = 1
                    paragraph.alignment = PP_ALIGN.LEFT
                    first_paragraph = False

    def _add_images(self, slide, images: Sequence[ImageInfo], top_inches: float) -> None:
        margin = Inches(0.9)
        available_width = Inches(self.width_in) - margin * 2
        top = Inches(top_inches)
        for src, caption in images:
            blob = self._decode_data_uri(src)
            if not blob:
                logger.warning({"message": "Skip image with invalid data URI."})
                continue
            picture = slide.shapes.add_picture(BytesIO(blob), margin, top, width=available_width)
            top = picture.top + picture.height + Inches(0.25)
            if caption:
                box = slide.shapes.add_textbox(margin, top, available_width, Inches(0.6))
                frame = box.text_frame
                frame.clear()
                paragraph = frame.paragraphs[0]
                paragraph.text = caption
                paragraph.font.size = Pt(16)
                paragraph.alignment = PP_ALIGN.LEFT
                top = box.top + box.height + Inches(0.3)

    # ------------------------------------------------------------------
    # Utilities
    # ------------------------------------------------------------------
    def _select_heading(self, slide_data: SlideData) -> str | None:
        return slide_data.headings[0][1] if slide_data.headings else None

    def _find_paragraph(self, slide_data: SlideData, class_name: str) -> str | None:
        for classes, text in slide_data.paragraphs:
            tokens = set(classes.split()) if classes else set()
            if class_name in tokens:
                return text
        return None

    def _iter_paragraphs(
        self,
        slide_data: SlideData,
        exclude_classes: set[str] | None = None,
    ) -> List[str]:
        results: list[str] = []
        for classes, text in slide_data.paragraphs:
            tokens = set(classes.split()) if classes else set()
            if exclude_classes and tokens.intersection(exclude_classes):
                continue
            results.append(text)
        return results

    def _decode_data_uri(self, uri: str | None) -> bytes | None:
        if not uri or not uri.startswith("data:"):
            return None
        try:
            _, data = uri.split(",", 1)
            return base64.b64decode(data)
        except Exception as exc:  # pragma: no cover - defensive
            logger.error(
                {
                    "message": "Failed to decode data URI.",
                    "status": "problem",
                    "error_message": str(exc),
                }
            )
            return None

