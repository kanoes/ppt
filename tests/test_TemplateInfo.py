import logging
from datetime import datetime
from pptx import Presentation

# ログの設定
logging.basicConfig(level=logging.INFO)


def test_slide_info(slide_number: int):
    """
    幻灯片情報を取得して出力するテスト関数。
    """
    try:
        # プレゼンテーションファイルを読み込む
        prs = Presentation("resources/smbc_template_new.pptx")

        # 指定されたスライドを取得
        slide = prs.slides[slide_number]

        # 現在時刻を取得し、フォーマット
        timestamp = datetime.now().strftime("%Y/%m/%d %H:%M")

        # テスト結果のデザイン
        border = "★" * 20
        section_separator = "✦" * 30

        # 情報を収集する
        slide_info = f"\n\n\n\n{border}🔷 テスト開始🔷{border}\n\n"
        slide_info += f"🕒 テスト時間: {timestamp}\n"
        slide_info += f"{section_separator}\n"
        slide_info += f"📌 テスト入力: スライド番号 {slide_number}\n"
        slide_info += f"{section_separator}\n"
        slide_info += "📄 テスト出力:\n"

        for i, shape in enumerate(slide.shapes):
            shape_info = f"\nシェイプ番号: {i}\n"
            shape_info += f"タイプ: {shape.shape_type}\n"

            if shape.has_text_frame:
                shape_info += f"内容: {shape.text}\n"

            # シェイプの位置（左上の座標）とサイズを出力
            shape_info += f"位置: (左: {shape.left}, 上: {shape.top})\n"
            shape_info += f"サイズ: (幅: {shape.width}, 高さ: {shape.height})\n"
            slide_info += shape_info

        placeholders_info = []
        for placeholder in slide.placeholders:
            placeholders_info.append(
                {
                    "インデックス": placeholder.placeholder_format.idx,
                    "タイプ": placeholder.placeholder_format.type,
                    "名前": placeholder.name,
                }
            )

        slide_info += "\nプレースホルダー情報:\n"
        for placeholder in placeholders_info:
            slide_info += f"インデックス: {placeholder['インデックス']}, タイプ: {placeholder['タイプ']}, 名前: {placeholder['名前']}\n"

        slide_info += f"{section_separator}\n"
        slide_info += f"\n{border}🔷 テスト完了🔷{border}\n\n\n\n"

        # 結果を命令行に出力
        print(slide_info)
    except Exception as e:
        # エラーが発生した場合、エラーメッセージを出力
        print("スライド情報取得中にエラーが発生しました:", e)


# ==============================
# テストデータ
# ==============================
class TestData:
    """
    テストデータを格納するクラス。
    """

    # ==============テストデータ================
    SLIDE_NUMBER = 22  # スライド番号


# テストの実行
test_slide_info(
    TestData.SLIDE_NUMBER,
)
